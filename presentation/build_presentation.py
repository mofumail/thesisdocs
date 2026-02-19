from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


REFS = {
    1: "M. Abufadda and K. Mansour, \"A Survey of Synthetic Data Generation for Machine Learning,\" ACIT, 2021. https://doi.org/10.1109/ACIT53391.2021.9677302",
    3: "J. Bobadilla et al., \"Creating synthetic datasets for collaborative filtering recommender systems using GANs,\" Knowledge-Based Systems, 2023. https://doi.org/10.1016/j.knosys.2023.111016",
    4: "D. Bountouridis et al., \"SIREN: A Simulation Framework for Understanding the Effects of Recommender Systems in Online News Environments,\" FAT*, 2019. https://doi.org/10.1145/3287560.3287583",
    7: "X. Chen et al., \"Generative Adversarial Reward Learning for Generalized Behavior Tendency Inference,\" arXiv:2105.00822, 2021. http://arxiv.org/abs/2105.00822",
    11: "A. Creswell et al., \"Generative Adversarial Networks: An Overview,\" IEEE Signal Processing Magazine, 2018. https://doi.org/10.1109/MSP.2017.2765202",
    13: "M. del Carmen Rodríguez-Hernández et al., \"DataGenCARS: A generator of synthetic data for the evaluation of context-aware recommendation systems,\" Pervasive and Mobile Computing, 2017. https://doi.org/10.1016/j.pmcj.2016.09.020",
    16: "A. Figueira and B. Vaz, \"Survey on Synthetic Data Generation, Evaluation Methods and GANs,\" Mathematics, 2022. https://doi.org/10.3390/math10152733",
    18: "M. Jakomin, T. Curk, and Z. Bosnić, \"Generating inter-dependent data streams for recommender systems,\" Simulation Modelling Practice and Theory, 2018. https://doi.org/10.1016/j.simpat.2018.07.013",
    19: "J. Jordon et al., \"Synthetic Data—What, why and how?\" arXiv:2205.03257, 2022. http://arxiv.org/abs/2205.03257",
    21: "Y. Lu et al., \"Machine Learning for Synthetic Data Generation: A Review,\" arXiv:2302.04062, 2024. http://arxiv.org/abs/2302.04062",
    22: "E. Lucherini et al., \"T-RECS: A Simulation Tool to Study the Societal Impact of Recommender Systems,\" arXiv:2107.08959, 2021. http://arxiv.org/abs/2107.08959",
    24: "J. McInerney et al., \"Accordion: A Trainable Simulator for Long-Term Interactive Systems,\" RecSys, 2021. https://doi.org/10.1145/3460231.3474259",
    25: "E. Stavinova et al., \"Synthetic Data-Based Simulators for Recommender Systems: A Survey,\" arXiv:2206.11338, 2022. http://arxiv.org/abs/2206.11338",
    27: "X. Zhao et al., \"UserSim: User Simulation via Supervised Generative Adversarial Network,\" WWW, 2021. https://doi.org/10.1145/3442381.345012",
}


def add_reference_box(slide_obj, refs) -> None:
    if not refs:
        return

    refs_text = " | ".join(f"[{r}] {REFS[r]}" for r in refs if r in REFS)
    if not refs_text:
        return

    # Bottom-left, small-font citation block
    box = slide_obj.shapes.add_textbox(Inches(0.45), Inches(6.78), Inches(12.1), Inches(0.58))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = refs_text
    p.font.size = Pt(8.4)
    p.font.color.rgb = RGBColor(95, 95, 95)


def build_deck(slides, output_pptx: Path) -> None:
    prs = Presentation()

    # 16:9 layout for typical projector use
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide in enumerate(slides):
        if i == 0:
            s = prs.slides.add_slide(prs.slide_layouts[0])
            title = s.shapes.title
            subtitle = s.placeholders[1]

            title.text = slide["title"]
            subtitle.text = slide["subtitle"]

            title_tf = title.text_frame
            title_tf.paragraphs[0].font.size = Pt(40)
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = RGBColor(20, 33, 61)

            sub_tf = subtitle.text_frame
            sub_tf.paragraphs[0].font.size = Pt(21)
            sub_tf.paragraphs[0].font.color.rgb = RGBColor(70, 70, 70)
        else:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            title = s.shapes.title
            body = s.shapes.placeholders[1]

            title.text = slide["title"]
            title_tf = title.text_frame
            title_tf.paragraphs[0].font.size = Pt(34)
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = RGBColor(20, 33, 61)

            tf = body.text_frame
            tf.clear()

            bullets = slide.get("bullets", [])
            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(23)
                p.font.color.rgb = RGBColor(35, 35, 35)

            # add speaker notes
            notes = slide.get("notes", "")
            if notes:
                notes_frame = s.notes_slide.notes_text_frame
                notes_frame.text = notes

            add_reference_box(s, slide.get("refs", []))

        if i == 0:
            add_reference_box(s, slide.get("refs", []))

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)


def write_markdown(slides, output_md: Path) -> None:
    lines = []
    lines.append("# Workshop Update Slides (NL)")

    for i, slide in enumerate(slides, start=1):
        lines.append(f"## Slide {i}: {slide['title']}")
        if i == 1:
            lines.append(slide["subtitle"])
        else:
            for b in slide.get("bullets", []):
                lines.append(f"- {b}")

        notes = slide.get("notes", "").strip()
        if notes:
            lines.append("**Spreeknotities:**")
            lines.append(notes)

        refs = slide.get("refs", [])
        if refs:
            lines.append("**Referenties op deze slide:**")
            for r in refs:
                if r in REFS:
                    lines.append(f"- [{r}] {REFS[r]}")
        lines.append("")

    output_md.parent.mkdir(parents=True, exist_ok=True)
    output_md.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    slides = [
        {
            "title": "Afstudeerupdate: Interaction Generator voor Recommender Systems",
            "subtitle": "Expert Workshop PvA tips + Projectvoorstel\n19 februari 2026\nBeau",
            "notes": "Doel in 1 zin: stand van zaken en voorgestelde onderzoeksrichting tonen.\nTiming: ~30 sec.",
            "refs": [],
        },
        {
            "title": "1) Korte context van het onderzoeksgebied",
            "bullets": [
                "Recommender systems worden meestal geëvalueerd met statische data, terwijl gedrag dynamisch is.",
                "Voor realistische evaluatie is synthetische interactiedata nodig (sessies, acties, timing).",
                "In onze setting: component (1) arrivals bestaat al; component (2) interaction generation moet verbeterd worden.",
            ],
            "notes": "Noem kort waarom dit relevant is voor robustere RS-evaluatie en training. Timing: ~60 sec.",
            "refs": [],
        },
        {
            "title": "2) Synthetic data landscape -> focus op RS",
            "bullets": [
                "Breed synthetic-data kader: surveys [1], [19], [21].",
                "Binnen RS: simulatie als aparte subliteratuur met eigen evaluatieproblemen.",
                "Doel voor deze thesis: interacties niet alleen voorspellen, maar volledige logs genereren.",
            ],
            "notes": "Houd dit hoog-over; details komen op volgende slides. Timing: ~60 sec.",
            "refs": [1, 19, 21],
        },
        {
            "title": "3) Bestaande RS-simulatie aanpakken",
            "bullets": [
                "Anchor: Stavinova survey [25] over synthetic-data simulators voor RS.",
                "Frameworks/simulators: SIREN [4], T-RECS [22], Accordion [24].",
                "Generators: DataGenCARS [13], Jakomin [18], plus GAN-richtingen.",
                "Conclusie: veel methodes, maar geen duidelijke convergentie op één dominante aanpak.",
            ],
            "notes": "Laat zien dat je het veld kent en waar de variatie zit. Timing: ~75 sec.",
            "refs": [25, 4, 22, 24, 13, 18],
        },
        {
            "title": "4) Literatuurkloof (kern)",
            "bullets": [
                "RS-simulatie literatuur en sequential recommendation literatuur zijn parallel ontwikkeld.",
                "Simulatiepapers genereren data, maar gebruiken vrijwel geen transformer-gebaseerde sequentiemodellen.",
                "Sequential recommendation papers (SASRec/BERT4Rec/BST) tonen sterke sequence-modellering, maar focussen op predictie i.p.v. generatie.",
                "Kloof: de modellen die gebruikerssequenties het best begrijpen, worden niet ingezet om die interacties te simuleren.",
            ],
            "notes": "Gebruik deze slide als centrale probleemdefinitie. Timing: ~75 sec.",
            "refs": [25],
        },
        {
            "title": "5) Eerste draft onderzoeksvraag",
            "bullets": [
                "Hoe en in welke mate kan een autoregressieve transformer realistische synthetische gebruikersinteractielogs genereren voor training en evaluatie van recommender systems?",
                "Subvraag: presteert dit beter dan regel/statistische en GAN-gebaseerde alternatieven?",
                "Werkhypothese: betere sequentiële fideliteit + betere downstream utility.",
            ],
            "notes": "Noem dat dit een eerste formulering is en feedback welkom is. Timing: ~60 sec.",
            "refs": [],
        },
        {
            "title": "6) Voorgestelde methode: multi-head autoregressieve transformer",
            "bullets": [
                "Te genereren tuple per tijdstap: (user_id, session_id, action_type, item_id, Δt, context_t).",
                "Head A: p(action_t | history)",
                "Head B: p(item_t | action_t, history)",
                "Head C: p(Δt_t | action_t, item_t, history)",
                "Voordeel: causal sequence generation + constraints (bv. geen purchase zonder eerdere view/cart).",
            ],
            "notes": "Dit is je methodeslide: maak het concreet en controleerbaar. Timing: ~90 sec.",
            "refs": [],
        },
        {
            "title": "7) Waarom AR-transformer i.p.v. GAN als hoofdroute?",
            "bullets": [
                "Discrete outputs: interactiedata is token-achtig (action/item), AR-modellen zijn hiervoor natuurlijk.",
                "Stabiliteit/diversiteit: minder gevoelig voor mode collapse dan GAN-training in deze context [11], [16].",
                "Controllability: prefix-conditioning en expliciete validiteitsregels tijdens generatie.",
                "Evaluatie: expliciete likelihood per stap maakt vergelijking en ablatie helderder.",
            ],
            "notes": "Koppel expliciet aan je gekozen probleemtype (structured event logs). Timing: ~75 sec.",
            "refs": [11, 16],
        },
        {
            "title": "8) Labomgeving + eerste indruk + next steps",
            "bullets": [
                "Labomgeving/dagelijkse routine: [kort mondeling toe te lichten].",
                "Eerste indruk: sterke basis in literatuur en duidelijke modelrichting.",
                "Volgende stappen: baselinevergelijking, fidelity-metrics, downstream RS-evaluatie.",
                "Feedbackvraag: welke evaluatiemetrieken zijn voor jullie essentieel om de simulator als overtuigend te zien?",
            ],
            "notes": "Hou dit persoonlijk en kort; sluit af met concrete feedbackvraag. Timing: ~45 sec.",
            "refs": [],
        },
    ]

    out_dir = Path(__file__).resolve().parent
    output_pptx = out_dir / "workshop_update_beau_2026-02-19.pptx"
    output_md = out_dir / "workshop_update_beau_2026-02-19.md"

    build_deck(slides, output_pptx)
    write_markdown(slides, output_md)

    print(f"Wrote: {output_pptx}")
    print(f"Wrote: {output_md}")


if __name__ == "__main__":
    main()
